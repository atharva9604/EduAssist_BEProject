"""
Document Processor - Extracts text from various document formats
Supports: PDF, PPT/PPTX, DOC/DOCX
"""
from typing import Union
from io import BytesIO
from pathlib import Path

try:
    from pypdf import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False


class DocumentProcessor:
    """Extracts text content from various document formats"""
    
    def __init__(self):
        if not PDF_AVAILABLE:
            raise ImportError("pypdf is required for PDF processing. Install with: pip install pypdf")
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PPT processing. Install with: pip install python-pptx")
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx is required for DOC processing. Install with: pip install python-docx")
    
    def extract_text(self, file_bytes: bytes, file_type: str) -> str:
        """
        Extract text from document bytes
        
        Args:
            file_bytes: Raw file content as bytes
            file_type: Type of file ('pdf', 'pptx', 'docx')
            
        Returns:
            Extracted text content
        """
        file_type = file_type.lower()
        
        if file_type == 'pdf':
            return self._extract_from_pdf(file_bytes)
        elif file_type in ('ppt', 'pptx'):
            return self._extract_from_pptx(file_bytes)
        elif file_type in ('doc', 'docx'):
            return self._extract_from_docx(file_bytes)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    def _extract_from_pdf(self, file_bytes: bytes) -> str:
        """Extract text from PDF file"""
        if not PDF_AVAILABLE:
            raise ImportError("pypdf is not available")
        
        try:
            pdf_file = BytesIO(file_bytes)
            reader = PdfReader(pdf_file)
            
            # Check if PDF is encrypted
            if reader.is_encrypted:
                try:
                    reader.decrypt("")  # Try empty password
                except Exception:
                    raise ValueError("PDF is encrypted/password protected. Please provide an unencrypted PDF.")
            
            full_text = ""
            for page in reader.pages:
                try:
                    page_text = page.extract_text()
                    if page_text:
                        full_text += page_text + "\n"
                except Exception as e:
                    print(f"Warning: Could not extract text from a page: {e}")
                    continue
            
            return full_text.strip()
        except Exception as e:
            raise ValueError(f"Failed to extract text from PDF: {str(e)}")
    
    def _extract_from_pptx(self, file_bytes: bytes) -> str:
        """Extract text from PowerPoint file"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not available")
        
        try:
            pptx_file = BytesIO(file_bytes)
            prs = Presentation(pptx_file)
            
            full_text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"--- Slide {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                    # Also check for text in tables
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    slide_text += cell.text + " "
                            slide_text += "\n"
                
                full_text += slide_text + "\n"
            
            return full_text.strip()
        except Exception as e:
            raise ValueError(f"Failed to extract text from PowerPoint: {str(e)}")
    
    def _extract_from_docx(self, file_bytes: bytes) -> str:
        """Extract text from Word document"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx is not available")
        
        try:
            docx_file = BytesIO(file_bytes)
            doc = Document(docx_file)
            
            full_text = ""
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text += paragraph.text + "\n"
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        full_text += " | ".join(row_text) + "\n"
                full_text += "\n"
            
            return full_text.strip()
        except Exception as e:
            raise ValueError(f"Failed to extract text from Word document: {str(e)}")
