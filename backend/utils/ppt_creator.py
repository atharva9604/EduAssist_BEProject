import os
import json
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

class PPTCreator:
    """Creates PowerPoint presentations from structured content data"""
    
    def __init__(self, storage_dir: str = "storage/presentations"):
        self.storage_dir = Path(storage_dir)
        self.storage_dir.mkdir(parents=True, exist_ok=True)
    
    def create_presentation(self, content_data: Dict, filename: str = None) -> str:
        """Create a PowerPoint presentation from structured content data"""
        try:
            # Create new presentation
            prs = Presentation()
            
            # Set slide size to widescreen (16:9)
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Generate filename if not provided
            if not filename:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                safe_title = "".join(c for c in content_data.get("presentation_title", "Presentation") if c.isalnum() or c in (' ', '-', '_')).rstrip()
                filename = f"{safe_title}_{timestamp}.pptx"
            
            # Create slides based on content data
            slides_data = content_data.get("slides", [])
            
            for slide_data in slides_data:
                slide_type = slide_data.get("slide_type", "content")
                
                if slide_type == "title":
                    self._create_title_slide(prs, slide_data, content_data)
                elif slide_type == "summary":
                    self._create_summary_slide(prs, slide_data)
                else:
                    self._create_content_slide(prs, slide_data)
            
            # Save presentation
            file_path = self.storage_dir / filename
            prs.save(str(file_path))
            
            return str(file_path)
            
        except Exception as e:
            print(f"Error creating presentation: {e}")
            raise e
    
    def _create_title_slide(self, prs: Presentation, slide_data: Dict, content_data: Dict):
        """Create a title slide"""
        # Use blank layout for custom title slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = content_data.get("presentation_title", slide_data.get("title", "Presentation"))
        
        # Format title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_run = title_paragraph.runs[0]
        title_run.font.size = Pt(44)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(31, 73, 125)  # Dark blue
        
        # Add subtitle if available
        subtitle = content_data.get("presentation_subtitle", "")
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            
            subtitle_paragraph = subtitle_frame.paragraphs[0]
            subtitle_paragraph.alignment = PP_ALIGN.CENTER
            subtitle_run = subtitle_paragraph.runs[0]
            subtitle_run.font.size = Pt(24)
            subtitle_run.font.color.rgb = RGBColor(68, 114, 196)  # Lighter blue
        
        # Add speaker notes if available
        speaker_notes = slide_data.get("speaker_notes", "")
        if speaker_notes:
            try:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.text = speaker_notes
                for paragraph in notes_text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
            except Exception as e:
                print(f"Warning: Could not add speaker notes to title slide: {e}")
        
        # Add decorative element
        self._add_decorative_shape(slide)
    
    def _create_content_slide(self, prs: Presentation, slide_data: Dict):
        """Create a content slide with title and bullet points"""
        # Use content layout
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "Slide Title")
        
        # Format title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(31, 73, 125)
        
        # Set content
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()
        
        # Add bullet points
        content_list = slide_data.get("content", [])
        for i, point in enumerate(content_list):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(68, 68, 68)

        # Add image if provided
        # MODE 5 IMAGE LOCK: Skip slides that explicitly forbid images
        slide_num = slide_data.get('slide_number', 0)
        if slide_data.get("_mode_5_image_mode") == "NONE" or slide_data.get("_no_image"):
            print(f"🚫 MODE 5 IMAGE LOCK: No image added to slide {slide_num} (No image declared)")
            # Remove any image_path that might have been set incorrectly
            if slide_data.get("image_path"):
                del slide_data["image_path"]
        else:
            image_path = slide_data.get("image_path")
            if image_path:
                # Convert to absolute path if relative
                abs_image_path = os.path.abspath(image_path) if not os.path.isabs(image_path) else image_path
                if os.path.exists(abs_image_path):
                    try:
                        print(f"📷 Adding image to slide {slide_data.get('slide_number')}: {abs_image_path}")
                        slide.shapes.add_picture(
                            abs_image_path,
                            left=Inches(8.0),
                            top=Inches(2.0),
                            width=Inches(4.5)
                        )
                        print(f"✓ Image added successfully to slide {slide_data.get('slide_number')}")
                    except Exception as e:
                        print(f"❌ Warning: Could not add image {abs_image_path}: {e}")
                        import traceback
                        traceback.print_exc()
                else:
                    print(f"⚠ Image path does not exist: {abs_image_path}")
                    print(f"  Current directory: {os.getcwd()}")
                    print(f"  Original path: {image_path}")

        # Add speaker notes if available
        speaker_notes = slide_data.get("speaker_notes", "")
        if speaker_notes:
            try:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.text = speaker_notes
                # Format speaker notes
                for paragraph in notes_text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = RGBColor(0, 0, 0)
            except Exception as e:
                print(f"Warning: Could not add speaker notes: {e}")
        
        # Visual suggestions removed - images not implemented yet
    
    def _create_summary_slide(self, prs: Presentation, slide_data: Dict):
        """Create a summary/conclusion slide"""
        # Use content layout
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "Summary")
        
        # Format title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(36)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(192, 0, 0)  # Red for emphasis
        
        # Set content
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()
        
        # Add summary points
        content_list = slide_data.get("content", [])
        for i, point in enumerate(content_list):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(31, 73, 125)
        
        # Add speaker notes if available
        speaker_notes = slide_data.get("speaker_notes", "")
        if speaker_notes:
            try:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.text = speaker_notes
                for paragraph in notes_text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = RGBColor(0, 0, 0)
            except Exception as e:
                print(f"Warning: Could not add speaker notes to summary slide: {e}")
    
    def _add_decorative_shape(self, slide):
        """Add a decorative shape to the slide"""
        try:
            # Add a rectangle at the bottom
            left = Inches(1)
            top = Inches(6.5)
            width = Inches(11.33)
            height = Inches(0.3)
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(68, 114, 196)
            shape.line.fill.background()
        except Exception as e:
            print(f"Error adding decorative shape: {e}")
    
    def _add_visual_note(self, slide, suggestion: str):
        """Add a visual suggestion as a note"""
        try:
            # Add a text box for visual suggestion
            left = Inches(8)
            top = Inches(5.5)
            width = Inches(4.5)
            height = Inches(1.5)
            
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.text = f"💡 Visual Suggestion:\n{suggestion}"
            
            # Format the note
            paragraph = text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.italic = True
            paragraph.font.color.rgb = RGBColor(128, 128, 128)
            
            # Add background color
            text_box.fill.solid()
            text_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
            
        except Exception as e:
            print(f"Error adding visual note: {e}")
    
    def get_presentation_info(self, file_path: str) -> Dict:
        """Get information about a presentation file"""
        try:
            prs = Presentation(file_path)
            return {
                "file_path": file_path,
                "filename": os.path.basename(file_path),
                "num_slides": len(prs.slides),
                "created_at": datetime.fromtimestamp(os.path.getctime(file_path)).isoformat(),
                "file_size": os.path.getsize(file_path)
            }
        except Exception as e:
            return {"error": str(e)}
    
    def list_presentations(self) -> List[Dict]:
        """List all presentations in the storage directory"""
        presentations = []
        try:
            for file_path in self.storage_dir.glob("*.pptx"):
                info = self.get_presentation_info(str(file_path))
                if "error" not in info:
                    presentations.append(info)
        except Exception as e:
            print(f"Error listing presentations: {e}")
        
        return sorted(presentations, key=lambda x: x["created_at"], reverse=True)

